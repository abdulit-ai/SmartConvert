import streamlit as st
from PIL import Image
import pytesseract
import pdfplumber
from docx import Document
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import io

# --------------------------------------------------
# PAGE CONFIG
# --------------------------------------------------
st.set_page_config(page_title="SmartConvert", layout="centered")

st.title("SmartConvert")
st.caption("Convert Anything. Edit Everything.")

# --------------------------------------------------
# SIDEBAR NAVIGATION
# --------------------------------------------------
page = st.sidebar.selectbox(
    "Choose Conversion Type",
    [
        "Image to Text",
        "Image to PDF",
        "Image to Word",
        "PDF to Text",
        "PDF to Word",
        "PDF to Excel",
        "PDF to Slides",
    ],
)

# ==================================================
# IMAGE TO TEXT
# ==================================================
if page == "Image to Text":
    st.header("Image to Text")

    uploaded = st.file_uploader("Upload Image", type=["png", "jpg", "jpeg"])

    if uploaded:
        image = Image.open(uploaded)
        st.image(image)

        with st.spinner("Extracting text..."):
            text = pytesseract.image_to_string(image)

        st.text_area("Extracted Text", text, height=200)

        st.download_button(
            "Download as TXT",
            text,
            file_name="extracted.txt",
            mime="text/plain",
        )

# ==================================================
# IMAGE TO PDF
# ==================================================
elif page == "Image to PDF":
    st.header("Image to PDF")

    uploaded = st.file_uploader("Upload Image", type=["png", "jpg", "jpeg"])

    if uploaded:
        image = Image.open(uploaded)

        pdf_bytes = io.BytesIO()
        image.convert("RGB").save(pdf_bytes, format="PDF")
        pdf_bytes.seek(0)

        st.download_button(
            "Download PDF",
            pdf_bytes,
            file_name="converted.pdf",
            mime="application/pdf",
        )

# ==================================================
# IMAGE TO WORD
# ==================================================
elif page == "Image to Word":
    st.header("Image to Word")

    uploaded = st.file_uploader("Upload Image", type=["png", "jpg", "jpeg"])

    if uploaded:
        image = Image.open(uploaded)
        text = pytesseract.image_to_string(image)

        doc = Document()
        doc.add_paragraph(text)

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        st.download_button(
            "Download Word File",
            buffer,
            file_name="converted.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

# ==================================================
# PDF TO TEXT
# ==================================================
elif page == "PDF to Text":
    st.header("PDF to Text")

    uploaded = st.file_uploader("Upload PDF", type=["pdf"])

    if uploaded:
        text = ""
        with pdfplumber.open(uploaded) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""

        st.text_area("Extracted Text", text, height=300)

        st.download_button(
            "Download TXT",
            text,
            file_name="pdf_text.txt",
            mime="text/plain",
        )

# ==================================================
# PDF TO WORD
# ==================================================
elif page == "PDF to Word":
    st.header("PDF to Word")

    uploaded = st.file_uploader("Upload PDF", type=["pdf"])

    if uploaded:
        text = ""
        with pdfplumber.open(uploaded) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""

        doc = Document()
        doc.add_paragraph(text)

        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)

        st.download_button(
            "Download Word",
            buffer,
            file_name="converted.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

# ==================================================
# PDF TO EXCEL
# ==================================================
elif page == "PDF to Excel":
    st.header("PDF to Excel")

    uploaded = st.file_uploader("Upload PDF", type=["pdf"])

    if uploaded:
        tables = []

        with pdfplumber.open(uploaded) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    tables.append(pd.DataFrame(table[1:], columns=table[0]))

        if tables:
            df = pd.concat(tables)

            buffer = io.BytesIO()
            df.to_excel(buffer, index=False)
            buffer.seek(0)

            st.download_button(
                "Download Excel",
                buffer,
                file_name="converted.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        else:
            st.warning("No tables detected in PDF.")

# ==================================================
# PDF TO SLIDES
# ==================================================
elif page == "PDF to Slides":
    st.header("PDF to Slides")

    uploaded = st.file_uploader("Upload PDF", type=["pdf"])

    if uploaded:
        prs = Presentation()

        with pdfplumber.open(uploaded) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Page Content"
                slide.placeholders[1].text = text

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        st.download_button(
            "Download Slides",
            buffer,
            file_name="converted.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
